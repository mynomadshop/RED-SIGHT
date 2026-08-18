from __future__ import annotations

import asyncio
import csv
import hashlib
import json
import os
import time
from pathlib import Path
from typing import Any

from PySide6.QtCore import Qt, QTimer
from PySide6.QtWidgets import (
    QFileDialog,
    QFrame,
    QHBoxLayout,
    QLabel,
    QLineEdit,
    QMessageBox,
    QPlainTextEdit,
    QPushButton,
    QSizePolicy,
    QTextEdit,
    QVBoxLayout,
    QWidget,
    QDockWidget,
)

from app.ui import action_palette_stage102 as s102


ROOT = Path(__file__).resolve().parents[2]
LOCALAPPDATA = Path(
    os.environ.get(
        "LOCALAPPDATA",
        str(Path.home() / "AppData" / "Local"),
    )
)
ATTACHMENT_STATE_DIR = (
    LOCALAPPDATA / "RedSight" / "memory" / "attachment-context"
)
ATTACHMENT_STATE_DIR.mkdir(parents=True, exist_ok=True)

MAX_ATTACHMENTS = 12
MAX_FILE_BYTES = 30 * 1024 * 1024
MAX_FILE_CHARS = 10500
MAX_ATTACHMENT_CONTEXT_CHARS = 30000

TEXT_EXTENSIONS = {
    ".txt", ".md", ".markdown", ".rst", ".py", ".pyi", ".js", ".jsx",
    ".ts", ".tsx", ".json", ".jsonl", ".yaml", ".yml", ".toml", ".ini",
    ".cfg", ".conf", ".csv", ".tsv", ".log", ".html", ".htm", ".css",
    ".scss", ".sql", ".xml", ".ps1", ".psm1", ".bat", ".cmd", ".sh",
    ".java", ".kt", ".kts", ".c", ".cc", ".cpp", ".h", ".hpp", ".cs",
    ".go", ".rs", ".rb", ".php", ".swift", ".gradle", ".properties",
    ".env.example", ".gitignore",
}
PDF_EXTENSIONS = {".pdf"}
WORD_EXTENSIONS = {".docx"}
EXCEL_EXTENSIONS = {".xlsx", ".xlsm"}
POWERPOINT_EXTENSIONS = {".pptx"}

SUPPORTED_EXTENSIONS = (
    TEXT_EXTENSIONS
    | PDF_EXTENSIONS
    | WORD_EXTENSIONS
    | EXCEL_EXTENSIONS
    | POWERPOINT_EXTENSIONS
)


def _input_text(widget) -> str:
    try:
        return s102._input_text(widget)
    except Exception:
        if isinstance(widget, QLineEdit):
            return widget.text()
        if isinstance(widget, (QTextEdit, QPlainTextEdit)):
            return widget.toPlainText()
        return ""


def _set_input_text(widget, text: str) -> None:
    try:
        s102._set_input_text(widget, text)
    except Exception:
        if isinstance(widget, QLineEdit):
            widget.setText(text)
        elif isinstance(widget, (QTextEdit, QPlainTextEdit)):
            widget.setPlainText(text)
        try:
            widget.setFocus()
        except Exception:
            pass


def _find_chat_input(window):
    widget = getattr(window, "_redsight_chat_input", None)
    if widget is not None:
        return widget

    try:
        widget = s102.s10.s91.s9.base._find_chat_input(window)
    except Exception:
        widget = None

    if widget is not None:
        window._redsight_chat_input = widget

    return widget


def _layout_item_contains_widget(item, target) -> bool:
    widget = item.widget()

    if widget is target:
        return True

    if widget is not None:
        try:
            if widget.isAncestorOf(target):
                return True
        except Exception:
            pass

    child_layout = item.layout()

    if child_layout is not None:
        return _layout_contains_widget(
            child_layout,
            target,
        )

    return False


def _layout_contains_widget(layout, target) -> bool:
    if layout is None:
        return False

    for index in range(layout.count()):
        if _layout_item_contains_widget(
            layout.itemAt(index),
            target,
        ):
            return True

    return False


def _direct_index_containing_widget(layout, target) -> int:
    if layout is None:
        return -1

    for index in range(layout.count()):
        if _layout_item_contains_widget(
            layout.itemAt(index),
            target,
        ):
            return index

    return -1


def _find_vertical_chat_stack(input_widget):
    """
    Find the QVBoxLayout that owns the main chat input row.

    This handles BOTH common Qt designs:
      1. the input sits inside a QWidget with its own QHBoxLayout;
      2. a QHBoxLayout containing the input is added directly to a QVBoxLayout.

    Stage 10.2 relied mostly on identifying the old transcript widget and
    therefore could fall back to a right-hand QDockWidget.  Stage 10.3 anchors
    placement to the query input itself.
    """
    child = input_widget

    for _depth in range(10):
        if child is None:
            break

        parent = child.parentWidget()

        if parent is None:
            break

        layout = parent.layout()

        if (
            isinstance(layout, QVBoxLayout)
            and _layout_contains_widget(
                layout,
                input_widget,
            )
        ):
            index = _direct_index_containing_widget(
                layout,
                input_widget,
            )

            if index >= 0:
                return parent, layout, index

        child = parent

    return None, None, -1


def _find_hbox_with_input(layout, input_widget):
    if layout is None:
        return None

    if isinstance(layout, QHBoxLayout):
        if _layout_contains_widget(
            layout,
            input_widget,
        ):
            return layout

    for index in range(layout.count()):
        child_layout = layout.itemAt(index).layout()

        if child_layout is not None:
            found = _find_hbox_with_input(
                child_layout,
                input_widget,
            )
            if found is not None:
                return found

        widget = layout.itemAt(index).widget()

        if widget is not None:
            child = widget.layout()
            if child is not None:
                found = _find_hbox_with_input(
                    child,
                    input_widget,
                )
                if found is not None:
                    return found

    return None


def _find_input_row_layout(input_widget):
    child = input_widget

    for _depth in range(8):
        if child is None:
            break

        parent = child.parentWidget()

        if parent is None:
            break

        layout = parent.layout()

        row = _find_hbox_with_input(
            layout,
            input_widget,
        )

        if row is not None:
            return parent, row

        child = parent

    return None, None


def _close_old_side_chat(window):
    dock = getattr(window, "_redsight_bubble_chat_dock", None)

    if dock is None:
        try:
            dock = window.findChild(QDockWidget, "RedSightBubbleChatDock")
        except Exception:
            dock = None

    if dock is not None:
        try:
            dock.hide()
            dock.setWidget(QWidget(dock))
            dock.deleteLater()
        except Exception:
            pass

    try:
        delattr(window, "_redsight_bubble_chat_dock")
    except Exception:
        pass


def _install_inline_bubble_view(window):
    existing = getattr(window, "_redsight_bubble_view", None)
    if existing is not None:
        return existing

    _close_old_side_chat(window)

    input_widget = _find_chat_input(window)
    bubble = s102.BubbleConversationView(window)
    bubble.setObjectName("RedSightInlineBubbleConversation")
    bubble.setMinimumHeight(260)
    bubble.setSizePolicy(
        QSizePolicy.Policy.Expanding,
        QSizePolicy.Policy.Expanding,
    )

    installed = False

    if input_widget is not None:
        _parent, stack, index = _find_vertical_chat_stack(input_widget)

        if stack is not None and index >= 0:
            try:
                stack.insertWidget(index, bubble, 1)
                window._redsight_inline_chat_stack = stack
                window._redsight_inline_chat_index = index
                installed = True
            except Exception as exc:
                s102._log_error("stage103-inline-chat-insert", exc)

    # Secondary path: replace the legacy transcript in-place if the current
    # Command Center layout does not expose a vertical input stack.
    if not installed:
        candidate = s102._find_existing_transcript(window)
        if candidate is not None:
            try:
                parent = candidate.parentWidget()
                layout = parent.layout() if parent is not None else None
                index = layout.indexOf(candidate) if layout is not None else -1
                if layout is not None and index >= 0:
                    layout.insertWidget(index, bubble, 1)
                    candidate.hide()
                    window._redsight_original_transcript_widget = candidate
                    installed = True
            except Exception as exc:
                s102._log_error("stage103-transcript-replace", exc)

    # If the normal input-stack path succeeded, hide a confidently identified
    # legacy plain-text transcript so only the red/blue inline bubbles remain.
    if installed and not hasattr(window, "_redsight_original_transcript_widget"):
        try:
            legacy = s102._find_existing_transcript(window)
            if legacy is not None and legacy is not bubble:
                legacy.hide()
                window._redsight_original_transcript_widget = legacy
        except Exception as exc:
            s102._log_error("stage103-hide-legacy-transcript", exc)

    # Never recreate the Stage 10.2 RIGHT-side conversation panel.
    # If a highly customized layout defeats both insertion methods, use a
    # bottom workspace dock so chat remains horizontally central.
    if not installed:
        dock = QDockWidget("REDSIGHT CHAT", window)
        dock.setObjectName("RedSightInlineChatFallbackDock")
        dock.setAllowedAreas(Qt.DockWidgetArea.BottomDockWidgetArea)
        dock.setWidget(bubble)
        dock.setMinimumHeight(340)
        window.addDockWidget(Qt.DockWidgetArea.BottomDockWidgetArea, dock)
        window._redsight_inline_chat_fallback_dock = dock

    window._redsight_bubble_view = bubble
    QTimer.singleShot(300, bubble.schedule_refresh)

    return bubble


def _session_state_path(session_id: str) -> Path:
    safe = "".join(
        ch for ch in str(session_id)
        if ch.isalnum() or ch in {"-", "_"}
    )[:180]
    return ATTACHMENT_STATE_DIR / (safe + ".json")


async def _active_session_id() -> str | None:
    try:
        result = await s102._request_async(
            "/memory/session/active",
            timeout=8,
        )
        session = result.get("session", {})
        sid = session.get("id")
        return str(sid) if sid else None
    except Exception:
        return None


def _load_session_file_paths(session_id: str | None) -> list[str]:
    if not session_id:
        return []

    path = _session_state_path(session_id)

    if not path.exists():
        return []

    try:
        data = json.loads(path.read_text(encoding="utf-8"))
        raw = data.get("paths", [])
        result = []
        for value in raw:
            candidate = Path(str(value))
            if candidate.exists() and candidate.is_file():
                result.append(str(candidate))
        return result[:MAX_ATTACHMENTS]
    except Exception:
        return []


def _save_session_file_paths(
    session_id: str | None,
    paths: list[str],
) -> None:
    if not session_id:
        return

    target = _session_state_path(session_id)

    clean = []
    seen = set()

    for value in paths:
        path = Path(str(value))
        key = os.path.normcase(os.path.abspath(str(path)))
        if key in seen:
            continue
        seen.add(key)
        if path.exists() and path.is_file():
            clean.append(str(path))
        if len(clean) >= MAX_ATTACHMENTS:
            break

    payload = {
        "session_id": session_id,
        "paths": clean,
        "updated_at": time.time(),
    }

    temp = target.with_suffix(".json.tmp")
    temp.write_text(
        json.dumps(payload, indent=2, ensure_ascii=False),
        encoding="utf-8",
    )
    os.replace(temp, target)


def _clear_session_file_paths(session_id: str | None) -> None:
    if not session_id:
        return

    try:
        _session_state_path(session_id).unlink(missing_ok=True)
    except Exception:
        pass


def _decode_text_bytes(data: bytes) -> str:
    for encoding in (
        "utf-8-sig",
        "utf-8",
        "utf-16",
        "cp1252",
        "latin-1",
    ):
        try:
            return data.decode(encoding)
        except Exception:
            continue

    return data.decode("utf-8", errors="replace")


def _extract_text_file(path: Path) -> str:
    # Read only a bounded prefix.  This prevents giant logs/code dumps from
    # exhausting the model context while still providing useful content.
    with path.open("rb") as handle:
        data = handle.read(2 * 1024 * 1024)

    if b"\x00" in data[:8192] and path.suffix.lower() not in {".txt", ".log"}:
        raise ValueError("File appears to be binary rather than text.")

    return _decode_text_bytes(data)


def _extract_pdf(path: Path) -> str:
    import pymupdf

    parts = []

    with pymupdf.open(path) as document:
        for page_number, page in enumerate(document):
            text = page.get_text("text") or ""
            if text.strip():
                parts.append(
                    f"\n--- PDF PAGE {page_number + 1} ---\n{text}"
                )
            if sum(len(part) for part in parts) >= MAX_FILE_CHARS:
                break

    return "\n".join(parts)


def _extract_docx(path: Path) -> str:
    import docx

    document = docx.Document(str(path))
    parts = []

    for paragraph in document.paragraphs:
        text = paragraph.text.strip()
        if text:
            parts.append(text)
        if sum(len(part) for part in parts) >= MAX_FILE_CHARS:
            break

    if sum(len(part) for part in parts) < MAX_FILE_CHARS:
        for table_index, table in enumerate(document.tables[:20], 1):
            parts.append(f"\n--- TABLE {table_index} ---")
            for row in table.rows[:100]:
                values = [cell.text.strip() for cell in row.cells]
                parts.append(" | ".join(values))
                if sum(len(part) for part in parts) >= MAX_FILE_CHARS:
                    break
            if sum(len(part) for part in parts) >= MAX_FILE_CHARS:
                break

    return "\n".join(parts)


def _extract_xlsx(path: Path) -> str:
    import openpyxl

    workbook = openpyxl.load_workbook(
        filename=str(path),
        read_only=True,
        data_only=True,
    )
    parts = []

    try:
        for sheet in workbook.worksheets[:10]:
            parts.append(f"\n--- SHEET: {sheet.title} ---")
            for row_number, row in enumerate(
                sheet.iter_rows(
                    min_row=1,
                    max_row=min(sheet.max_row or 1, 500),
                    values_only=True,
                ),
                1,
            ):
                values = []
                for value in row[:40]:
                    if value is None:
                        values.append("")
                    else:
                        values.append(str(value))
                if any(value for value in values):
                    parts.append(
                        str(row_number) + ": " + " | ".join(values)
                    )
                if sum(len(part) for part in parts) >= MAX_FILE_CHARS:
                    break
            if sum(len(part) for part in parts) >= MAX_FILE_CHARS:
                break
    finally:
        workbook.close()

    return "\n".join(parts)


def _extract_pptx(path: Path) -> str:
    from pptx import Presentation

    presentation = Presentation(str(path))
    parts = []

    for index, slide in enumerate(presentation.slides, 1):
        slide_parts = []

        for shape in slide.shapes:
            text = getattr(shape, "text", "")
            if isinstance(text, str) and text.strip():
                slide_parts.append(text.strip())

        if slide_parts:
            parts.append(
                f"\n--- SLIDE {index} ---\n"
                + "\n".join(slide_parts)
            )

        if sum(len(part) for part in parts) >= MAX_FILE_CHARS:
            break

    return "\n".join(parts)


def _sha256(path: Path) -> str:
    digest = hashlib.sha256()

    with path.open("rb") as handle:
        while True:
            chunk = handle.read(1024 * 1024)
            if not chunk:
                break
            digest.update(chunk)

    return digest.hexdigest()


def _extract_file_context(path_value: str) -> tuple[str, dict[str, Any]]:
    path = Path(path_value)

    if not path.exists() or not path.is_file():
        raise FileNotFoundError(str(path))

    stat = path.stat()

    if stat.st_size > MAX_FILE_BYTES:
        raise ValueError(
            f"File is {stat.st_size / (1024 * 1024):.1f} MB; "
            f"the direct-chat attachment limit is {MAX_FILE_BYTES / (1024 * 1024):.0f} MB."
        )

    suffix = path.suffix.lower()

    if suffix in PDF_EXTENSIONS:
        text = _extract_pdf(path)
        parser = "pymupdf"
    elif suffix in WORD_EXTENSIONS:
        text = _extract_docx(path)
        parser = "python-docx"
    elif suffix in EXCEL_EXTENSIONS:
        text = _extract_xlsx(path)
        parser = "openpyxl"
    elif suffix in POWERPOINT_EXTENSIONS:
        text = _extract_pptx(path)
        parser = "python-pptx"
    else:
        text = _extract_text_file(path)
        parser = "text"

    text = str(text).strip()

    if not text:
        raise ValueError(
            "No text could be extracted from this file. "
            "Image-only/scanned documents require a vision/OCR ingestion path."
        )

    truncated = len(text) > MAX_FILE_CHARS
    text = text[:MAX_FILE_CHARS]

    metadata = {
        "name": path.name,
        "path": str(path),
        "suffix": suffix,
        "size_bytes": int(stat.st_size),
        "modified": float(stat.st_mtime),
        "sha256": _sha256(path),
        "parser": parser,
        "truncated": truncated,
    }

    header = (
        "[REDSIGHT CHAT ATTACHMENT]\n"
        f"Name: {metadata['name']}\n"
        f"Source path: {metadata['path']}\n"
        f"Size: {metadata['size_bytes']} bytes\n"
        f"SHA256: {metadata['sha256']}\n"
        f"Parser: {metadata['parser']}\n"
        f"Truncated for direct model context: {metadata['truncated']}\n"
        "Content:\n"
    )

    return header + text, metadata


def _build_attachment_context(paths: list[str]) -> tuple[str, list[dict[str, Any]]]:
    sections = []
    metadata = []
    used = 0

    for value in paths[:MAX_ATTACHMENTS]:
        try:
            section, item = _extract_file_context(value)
        except Exception as exc:
            item = {
                "name": Path(value).name,
                "path": str(value),
                "error": type(exc).__name__ + ": " + str(exc),
            }
            metadata.append(item)
            section = (
                "[REDSIGHT CHAT ATTACHMENT ERROR]\n"
                f"Name: {Path(value).name}\n"
                f"Path: {value}\n"
                f"Error: {item['error']}\n"
            )
        else:
            metadata.append(item)

        remaining = MAX_ATTACHMENT_CONTEXT_CHARS - used
        if remaining <= 0:
            break

        clipped = section[:remaining]
        sections.append(clipped)
        used += len(clipped)

    context = "\n\n".join(sections).strip()

    if context:
        context = (
            "[ATTACHED FILE CONTEXT]\n"
            "Treat the following local file extracts as user-provided context. "
            "Do not claim you opened or modified files beyond the extraction shown. "
            "When answering, distinguish file contents from your inference.\n\n"
            + context
        )

    return context, metadata


def _attachment_names(paths: list[str]) -> str:
    names = [Path(value).name for value in paths]
    if len(names) <= 4:
        return ", ".join(names)
    return ", ".join(names[:4]) + f" +{len(names) - 4} more"


def _refresh_attachment_tray(window) -> None:
    paths = list(getattr(window, "_redsight_pending_attachments", []))
    label = getattr(window, "_redsight_attachment_label", None)
    clear_button = getattr(window, "_redsight_attachment_clear_button", None)

    if label is not None:
        if paths:
            label.setText(
                f"📎 {len(paths)} file(s) in this chat context  •  "
                + _attachment_names(paths)
            )
            label.setToolTip("\n".join(paths))
            label.show()
        else:
            label.setText("No files attached")
            label.setToolTip("")
            label.hide()

    if clear_button is not None:
        clear_button.setVisible(bool(paths))


async def _sync_attachment_session(window) -> str | None:
    session_id = await _active_session_id()
    previous = getattr(window, "_redsight_attachment_session_id", None)

    if session_id and session_id != previous:
        window._redsight_attachment_session_id = session_id
        window._redsight_pending_attachments = _load_session_file_paths(
            session_id
        )
        _refresh_attachment_tray(window)

    return session_id


def _select_attachments(window) -> None:
    current = list(getattr(window, "_redsight_pending_attachments", []))

    filters = (
        "RedSight context files "
        "(*.pdf *.docx *.xlsx *.xlsm *.pptx *.txt *.md *.markdown *.rst "
        "*.py *.js *.jsx *.ts *.tsx *.json *.jsonl *.yaml *.yml *.toml "
        "*.ini *.cfg *.conf *.csv *.tsv *.log *.html *.htm *.css *.sql "
        "*.xml *.ps1 *.psm1 *.bat *.cmd *.sh *.java *.kt *.c *.cpp *.h "
        "*.hpp *.cs *.go *.rs);;"
        "All files (*.*)"
    )

    selected, _filter = QFileDialog.getOpenFileNames(
        window,
        "Attach files to RedSight chat context",
        str(Path.home()),
        filters,
    )

    if not selected:
        return

    merged = []
    seen = set()

    for value in current + list(selected):
        path = Path(value)
        key = os.path.normcase(os.path.abspath(str(path)))

        if key in seen:
            continue
        seen.add(key)

        if not path.exists() or not path.is_file():
            continue

        if path.stat().st_size > MAX_FILE_BYTES:
            QMessageBox.warning(
                window,
                "RedSight attachment too large",
                (
                    f"{path.name} is larger than the "
                    f"{MAX_FILE_BYTES // (1024 * 1024)} MB direct-chat limit."
                ),
            )
            continue

        suffix = path.suffix.lower()

        # Unknown extensions are allowed only when they look like text.
        if suffix not in SUPPORTED_EXTENSIONS:
            try:
                probe = path.read_bytes()[:8192]
                if b"\x00" in probe:
                    QMessageBox.warning(
                        window,
                        "Unsupported RedSight attachment",
                        (
                            f"{path.name} appears to be a binary file and "
                            "cannot be injected into the text chat context."
                        ),
                    )
                    continue
            except Exception:
                continue

        merged.append(str(path))

        if len(merged) >= MAX_ATTACHMENTS:
            break

    window._redsight_pending_attachments = merged
    _refresh_attachment_tray(window)

    async def persist():
        sid = await _active_session_id()
        if sid:
            window._redsight_attachment_session_id = sid
            _save_session_file_paths(sid, merged)

    try:
        asyncio.create_task(persist())
    except RuntimeError:
        pass


def _clear_attachments(window) -> None:
    window._redsight_pending_attachments = []
    _refresh_attachment_tray(window)

    async def persist():
        sid = await _active_session_id()
        _clear_session_file_paths(sid)

    try:
        asyncio.create_task(persist())
    except RuntimeError:
        pass


def _install_attachment_controls(window):
    if hasattr(window, "_redsight_attach_button"):
        return window._redsight_attach_button

    input_widget = _find_chat_input(window)

    if input_widget is None:
        return None

    row_parent, row_layout = _find_input_row_layout(input_widget)
    _stack_parent, stack_layout, anchor_index = (
        _find_vertical_chat_stack(input_widget)
    )

    attach_button = QPushButton("📎 Attach")
    attach_button.setObjectName("RedSightAttachFilesButton")
    attach_button.setToolTip(
        "Attach local files as persistent context for the active RedSight chat."
    )
    attach_button.setStyleSheet(
        """
        QPushButton#RedSightAttachFilesButton {
            color:#EAF4FF;
            background:#124C88;
            border:1px solid #58A6F6;
            border-radius:9px;
            padding:7px 10px;
            font-weight:800;
        }
        QPushButton#RedSightAttachFilesButton:hover {
            background:#1763AE;
        }
        """
    )

    if row_layout is not None:
        input_index = row_layout.indexOf(input_widget)
        if input_index < 0:
            input_index = 0
        row_layout.insertWidget(input_index, attach_button)
    else:
        # If the input itself is a direct child of a vertical stack, put the
        # button immediately before it.  The common layout uses a horizontal
        # row, so this is only a compatibility path.
        parent = input_widget.parentWidget()
        layout = parent.layout() if parent is not None else None
        if layout is not None:
            index = layout.indexOf(input_widget)
            if index >= 0:
                layout.insertWidget(index, attach_button)

    tray = QFrame()
    tray.setObjectName("RedSightAttachmentTray")
    tray.setStyleSheet(
        """
        QFrame#RedSightAttachmentTray {
            background:#0B1723;
            border:1px solid #24577D;
            border-radius:8px;
        }
        QLabel {
            color:#A9D5FF;
            background:transparent;
            border:none;
            font-size:11px;
        }
        QPushButton {
            color:#FFFFFF;
            background:#5D1920;
            border:1px solid #B6444E;
            border-radius:7px;
            padding:4px 8px;
            font-weight:700;
        }
        """
    )

    tray_layout = QHBoxLayout(tray)
    tray_layout.setContentsMargins(8, 4, 6, 4)
    tray_layout.setSpacing(6)

    attachment_label = QLabel("No files attached")
    attachment_label.setWordWrap(False)

    clear_button = QPushButton("Clear Files")
    clear_button.setToolTip(
        "Remove attached file context from the active RedSight conversation."
    )

    tray_layout.addWidget(attachment_label, 1)
    tray_layout.addWidget(clear_button)

    # Put the attachment tray directly above the full query row.
    if (
        stack_layout is not None
        and anchor_index >= 0
    ):
        stack_layout.insertWidget(anchor_index, tray)
    elif row_parent is not None:
        # Compatibility fallback: insert immediately before the input row's
        # parent whenever its own parent exposes a vertical layout.
        parent = row_parent.parentWidget()
        layout = parent.layout() if parent is not None else None
        if isinstance(layout, QVBoxLayout):
            index = layout.indexOf(row_parent)
            if index >= 0:
                layout.insertWidget(index, tray)

    attach_button.clicked.connect(
        lambda: _select_attachments(window)
    )
    clear_button.clicked.connect(
        lambda: _clear_attachments(window)
    )

    window._redsight_attach_button = attach_button
    window._redsight_attachment_tray = tray
    window._redsight_attachment_label = attachment_label
    window._redsight_attachment_clear_button = clear_button
    window._redsight_pending_attachments = []

    _refresh_attachment_tray(window)

    async def initial_sync():
        await _sync_attachment_session(window)

    QTimer.singleShot(
        300,
        lambda: asyncio.create_task(initial_sync()),
    )

    return attach_button


def _display_user_message(
    original: str,
    paths: list[str],
) -> str:
    if not paths:
        return original

    return (
        original
        + "\n\n📎 Attached context: "
        + _attachment_names(paths)
    ).strip()


async def _send_dispatch(self, message):
    # Sync attached-file state with the active persistent conversation before
    # each turn.  This prevents files from leaking across /new sessions.
    await _sync_attachment_session(self)

    paths = list(
        getattr(
            self,
            "_redsight_pending_attachments",
            [],
        )
    )

    original = str(message).strip()

    if not original and paths:
        original = (
            "Review the attached file context and summarize the most relevant "
            "information, issues, and next actions."
        )

    if not original:
        return None

    if bool(getattr(self, "_redsight_stage103_inflight", False)):
        view = getattr(self, "_redsight_bubble_view", None)
        if view is not None:
            view.append_local(
                "assistant",
                "RedSight is still processing the previous request. "
                "Wait for it to finish before sending another action.",
            )
        return None

    self._redsight_stage103_inflight = True
    self._redsight_sounds_enabled = bool(
        getattr(
            self,
            "_redsight_sounds_enabled",
            True,
        )
    )

    view = getattr(self, "_redsight_bubble_view", None)
    indicator = getattr(self, "_redsight_processing91", None)

    display_original = _display_user_message(original, paths)

    if view is not None:
        view.append_local("user", display_original)

    s102._play_tone("send", self)

    if indicator is not None:
        try:
            indicator.begin(
                s102.s10.s91.description(original)
            )
        except Exception:
            try:
                indicator.begin("PROCESSING")
            except Exception:
                pass

    session_id: str | None = None
    effective = original
    attachment_context = ""
    attachment_metadata: list[dict[str, Any]] = []

    try:
        if s102._is_reset_request(original):
            await s102._new_session(self, original)
            self._redsight_pending_attachments = []
            self._redsight_attachment_session_id = await _active_session_id()
            _refresh_attachment_tray(self)
            s102._play_tone("receive", self)
            return "New conversation started."

        if original.lower().split(" ", 1)[0] == "/restart":
            return await s102._restart_platform(self, original)

        if original.lower().startswith("/errors"):
            try:
                text = s102.CHAT_LOG.read_text(
                    encoding="utf-8",
                    errors="replace",
                )
                assistant = (
                    "Recent RedSight UI diagnostics:\n\n"
                    + (
                        text[-12000:]
                        if text.strip()
                        else "No Stage 10.2/10.3 UI errors have been logged."
                    )
                )
            except Exception as exc:
                assistant = s102._error_text(
                    "read-diagnostics",
                    exc,
                )

            if view is not None:
                view.append_local("assistant", assistant)

            s102._play_tone("receive", self)
            return assistant

        if original.lower().startswith("/sessions"):
            sessions = await s102._request_async(
                "/memory/sessions",
                timeout=10,
            )
            assistant = s102._session_list_text(sessions)

            if view is not None:
                view.append_local("assistant", assistant)

            s102._play_tone("receive", self)
            return assistant

        if paths:
            attachment_context, attachment_metadata = await asyncio.to_thread(
                _build_attachment_context,
                paths,
            )

            if attachment_context:
                effective = (
                    original
                    + "\n\n"
                    + attachment_context
                )[:35500]

        auto_agent = bool(
            getattr(
                self,
                "_redsight_agent_mode",
                False,
            )
        )
        is_slash = original.startswith("/")

        # Plain chat + attachments stays on the deterministic chat path so the
        # actual extracted file text reaches LM Studio.  Explicit slash/agent
        # commands still use the action router.
        action_intent = (
            is_slash
            or (
                not paths
                and (
                    auto_agent
                    or s102.s10.looks_actionable_stage10(original)
                )
            )
        )

        if action_intent:
            command = (
                original
                if is_slash
                else "/agent " + original
            )

            if command.split(" ", 1)[0].lower() in {
                "/new",
                "/reset",
                "/clear",
            }:
                await s102._new_session(self, original)
                self._redsight_pending_attachments = []
                self._redsight_attachment_session_id = await _active_session_id()
                _refresh_attachment_tray(self)
                s102._play_tone("receive", self)
                return "New conversation started."

            result = await s102.s10.handle_stage10_slash(
                self,
                command,
            )

            if (
                command.split(" ", 1)[0].lower()
                in {"/help", "/actions"}
                and isinstance(result, dict)
            ):
                text = str(result.get("help", ""))
                additions = (
                    "\n📎 Attach Files"
                    "\n    Use the Attach button in the query bar to add "
                    "PDF/DOCX/XLSX/PPTX/text/code context."
                )
                if "📎 Attach Files" not in text:
                    result["help"] = text + additions

            direct = s102._direct_result(
                command,
                result,
            )

            if direct is not None:
                effective = s102._effective_action_message(
                    display_original,
                    command,
                    result,
                )
                session_id, _messages = await s102._build_messages(
                    display_original,
                    effective,
                )
                await s102._commit_turn(
                    display_original,
                    direct,
                    effective,
                    session_id,
                )

                if view is not None:
                    await view.refresh_async(force=True)

                s102._play_tone("receive", self)
                return direct

            effective = s102._effective_action_message(
                display_original,
                command,
                result,
            )

        # For a normal file-context turn, effective already contains the
        # bounded extracted text.
        session_id, messages = await s102._build_messages(
            display_original,
            effective,
        )

        assistant, _data = await s102._backend_chat(
            messages
        )

        await s102._commit_turn(
            display_original,
            assistant,
            effective,
            session_id,
        )

        # Persist the selected attachment paths per session.  Files remain
        # active chat context until Clear Files or /new is used.
        if paths and session_id:
            _save_session_file_paths(
                session_id,
                paths,
            )
            self._redsight_attachment_session_id = session_id

        if view is not None:
            await view.refresh_async(force=True)

        dock = getattr(
            self,
            "_redsight_conversation_memory_dock",
            None,
        )
        if dock is not None:
            QTimer.singleShot(100, dock.refresh)

        s102._play_tone("receive", self)
        return assistant

    except Exception as exc:
        s102._log_error(
            "stage103-chat-dispatch",
            exc,
            "User input: " + original[:3000],
        )
        assistant = s102._error_text(
            "Stage 10.3 chat/file-context pipeline",
            exc,
        )

        try:
            if session_id is None:
                built = await s102._request_async(
                    "/memory/build",
                    body={
                        "user_message": display_original,
                        "effective_message": effective,
                        "heritage_context": s102._heritage_context(
                            original
                        ),
                    },
                    timeout=10,
                )
                session_id = built.get("session_id")

            await s102._commit_turn(
                display_original,
                assistant,
                effective,
                session_id,
            )
        except Exception as commit_exc:
            s102._log_error(
                "stage103-error-turn-commit",
                commit_exc,
            )

        if view is not None:
            try:
                await view.refresh_async(force=True)
            except Exception:
                view.append_local(
                    "assistant",
                    assistant,
                )

        s102._play_tone("error", self)
        return assistant

    finally:
        if indicator is not None:
            try:
                indicator.finish()
            except Exception:
                pass

        self._redsight_stage103_inflight = False


def install_action_hooks(command_center_class):
    if getattr(
        command_center_class,
        "_redsight_stage103_installed",
        False,
    ):
        return

    # Stage 10.3 remains a single dispatcher: no old nested action-summary
    # wrapper is reinstalled.
    command_center_class._send_to_api = _send_dispatch
    command_center_class._redsight_stage103_installed = True


def attach_action_palette(window, project_root):
    # Force Stage 10.2's chat view installer to use the new inline location
    # BEFORE it constructs any conversation widget.
    s102._install_bubble_view = _install_inline_bubble_view

    palette = s102.attach_action_palette(
        window,
        project_root,
    )

    _install_attachment_controls(window)

    # Defensive cleanup in case a historical saved UI state reintroduced the
    # old right-side dock.
    _close_old_side_chat(window)

    return palette
